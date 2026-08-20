from __future__ import annotations

import ipaddress
import os
import re
import socket
import unicodedata
from pathlib import Path
from urllib.parse import urlparse

import requests


MAX_WEB_BYTES = 5 * 1024 * 1024


def normalize_pdf_text(text: str) -> str:
    """Normalize extracted PDF text while preserving its meaningful lines."""
    text = unicodedata.normalize("NFKC", text).replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.splitlines()]
    normalized = []
    for line in lines:
        if not line:
            if normalized and normalized[-1]:
                normalized.append("")
        elif normalized and normalized[-1].endswith("-") and line[:1].islower():
            normalized[-1] = normalized[-1][:-1] + line
        else:
            normalized.append(line)
    return "\n".join(normalized).strip()


def extract_pdf_title(text: str, fallback_title: str) -> str:
    """Use the bracketed PDF header as the document-wide title."""
    match = re.search(r"\[([^\]]{2,200})\]", text)
    return match.group(1).strip() if match else fallback_title


def _ocr_answer_line(line: str) -> str | None:
    """Normalize OCR variants of the FAQ answer marker (☞)."""
    match = re.match(r"^[☞>▷▶ㅠㅜ~]\s*(.*)$", line)
    if match:
        return match.group(1).strip()
    return None


def structure_pdf_entries(text: str, document_title: str) -> list[str]:
    """Group OCR lines into FAQ entries using their visual marker semantics."""
    entries = []
    current = None
    section = ""
    active = ""

    def finish_current():
        if not current or not current["question"]:
            return
        parts = [f"문서 제목: {document_title}"]
        if current["section"]:
            parts.append(f"분류: {current['section']}")
        parts.append(f"질문: {current['question']}")
        parts.append(f"답변: {' '.join(current['answer']) or '답변이 없습니다.'}")
        if current["notes"]:
            parts.append(f"기타 참고사항: {' '.join(current['notes'])}")
        entries.append("\n\n".join(parts))

    for line in normalize_pdf_text(text).splitlines():
        if not line:
            continue
        if line.startswith("[") and "]" in line:
            continue
        if line.startswith("*"):
            if current:
                current["notes"].append(line.lstrip("*").strip())
                active = "notes"
            continue
        answer = _ocr_answer_line(line)
        if answer is not None:
            if current:
                current["answer"].append(answer)
                active = "answer"
            continue
        if line.endswith("?"):
            finish_current()
            current = {"section": section, "question": line, "answer": [], "notes": []}
            active = "question"
            continue
        if current and active == "answer":
            current["answer"].append(line)
        elif current and active == "notes":
            current["notes"].append(line)
        elif not current:
            section = line

    finish_current()
    return entries or [f"문서 제목: {document_title}\n\n본문: {normalize_pdf_text(text)}"]


_COMPACT_QUESTION = re.compile(
    r"(?<!\d)\d{1,2}\s*(?=[\uac00-\ud7a3A-Za-z\[])[^?\n]{5,300}\?"
)
_NEXT_QUESTION = _COMPACT_QUESTION
_COMPACT_QUESTION_START = re.compile(
    r"(?<!\d)\d{1,2}(?:\s+\d{1,2})?\s*(?=[\uac00-\ud7a3A-Za-z\[])"
)
_COMPACT_TOP_LEVEL_QUESTION_START = re.compile(
    r"(?:^|(?<=[.!?]))\s*(\d{1,2}(?:\s+\d{1,2})?\s*(?=[\uac00-\ud7a3A-Za-z\[]))"
)


def structure_compact_faq_entries(text: str, document_title: str) -> list[str]:
    """Parse a text-layer FAQ where visual line breaks were flattened."""
    text = normalize_pdf_text(text).replace("\n", " ")
    markers = list(re.finditer("☞", text))
    entries = []
    previous_marker_end = 0

    for index, marker in enumerate(markers):
        question_region = text[previous_marker_end:marker.start()]
        question_starts = list(_COMPACT_QUESTION_START.finditer(question_region))
        if not question_starts:
            previous_marker_end = marker.end()
            continue
        top_level_question_starts = list(
            _COMPACT_TOP_LEVEL_QUESTION_START.finditer(question_region)
        )
        # The final numbered segment before an answer marker is the question.
        # Earlier numeric phrases ("1기", "3개" etc.) belong to the prior answer.
        # A page heading such as "2장 ..." can be glued to the first FAQ.
        # In that case use the final numbered start, which is the actual FAQ.
        if re.match(r"^\d{1,2}\s*\uc7a5", question_region) and len(question_starts) > 1:
            question_start = question_starts[-1].start()
        else:
            question_start = (
                top_level_question_starts[-1].start(1)
                if top_level_question_starts
                else question_starts[-1].start()
            )
        question = question_region[question_start:].strip()
        answer_end = markers[index + 1].start() if index + 1 < len(markers) else len(text)
        answer_region = text[marker.end():answer_end]
        # Text-layer extraction may join the next numbered question directly
        # to the previous answer (for example: "...answer. 4 next question?").
        # Keep that next question in its own FAQ entry instead of embedding it
        # in the previous answer.
        next_question_starts = list(_COMPACT_QUESTION_START.finditer(answer_region))
        if next_question_starts:
            next_top_level_starts = list(
                _COMPACT_TOP_LEVEL_QUESTION_START.finditer(answer_region)
            )
            next_question_start = (
                next_top_level_starts[-1].start(1)
                if next_top_level_starts
                else next_question_starts[-1].start()
            )
            answer_region = answer_region[:next_question_start]
        answer, separator, notes = answer_region.partition("*")
        parts = [
            f"문서 제목: {document_title}",
            f"질문: {question}",
            f"답변: {answer.strip() or '답변이 없습니다.'}",
        ]
        if separator and notes.strip():
            parts.append(f"기타 참고사항: {notes.strip()}")
        entries.append("\n\n".join(parts))
        previous_marker_end = marker.end()
    return entries


def structure_pdf_text(text: str, fallback_title: str) -> str:
    """Backward-compatible single-text representation for tests and callers."""
    return "\n\n---\n\n".join(structure_pdf_entries(text, fallback_title))


def _configure_tesseract(pytesseract_module) -> None:
    for candidate in (os.environ.get("TESSERACT_CMD"), r"C:\Program Files\Tesseract-OCR\tesseract.exe"):
        if candidate and Path(candidate).is_file():
            pytesseract_module.pytesseract.tesseract_cmd = candidate
            return
    raise ValueError("Tesseract OCR 엔진을 찾지 못했습니다.")


def _ocr_pdf_page(pdf, page_number: int, tessdata_dir: Path) -> str:
    import pymupdf
    import pytesseract
    from PIL import Image

    _configure_tesseract(pytesseract)
    if not all((tessdata_dir / f"{language}.traineddata").is_file() for language in ("kor", "eng")):
        raise ValueError("한국어·영어 OCR 데이터를 찾지 못했습니다.")
    os.environ["TESSDATA_PREFIX"] = str(tessdata_dir)
    page = pdf.load_page(page_number)
    pixmap = page.get_pixmap(
        matrix=pymupdf.Matrix(3, 3), colorspace=pymupdf.csRGB, alpha=False
    )
    image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
    return pytesseract.image_to_string(image, lang="kor+eng", config="--psm 6")


def validate_public_url(url: str) -> str:
    parsed = urlparse(url.strip())
    if parsed.scheme not in {"http", "https"} or not parsed.hostname:
        raise ValueError("http 또는 https URL만 사용할 수 있습니다.")
    if parsed.username or parsed.password:
        raise ValueError("인증정보가 포함된 URL은 사용할 수 없습니다.")
    try:
        addresses = {
            item[4][0] for item in socket.getaddrinfo(parsed.hostname, parsed.port or 443)
        }
    except socket.gaierror as error:
        raise ValueError("URL의 호스트를 확인할 수 없습니다.") from error
    for address in addresses:
        ip = ipaddress.ip_address(address)
        if not ip.is_global:
            raise ValueError("내부망·로컬 주소는 수집할 수 없습니다.")
    return parsed.geturl()


def load_url(url: str) -> list:
    from bs4 import BeautifulSoup
    from langchain_core.documents import Document

    url = validate_public_url(url)
    with requests.get(
        url,
        headers={"User-Agent": "EXAONE-Django-RAG/1.0"},
        timeout=(5, 20),
        stream=True,
        allow_redirects=False,
    ) as response:
        if 300 <= response.status_code < 400:
            raise ValueError("리디렉션 URL은 수집할 수 없습니다. 최종 URL을 입력하세요.")
        response.raise_for_status()
        content_type = response.headers.get("Content-Type", "").lower()
        if "text/html" not in content_type and "text/plain" not in content_type:
            raise ValueError("HTML 또는 텍스트 URL만 수집할 수 있습니다.")
        body = bytearray()
        for chunk in response.iter_content(64 * 1024):
            body.extend(chunk)
            if len(body) > MAX_WEB_BYTES:
                raise ValueError("URL 본문은 5MB 이하여야 합니다.")
    response.encoding = response.encoding or response.apparent_encoding or "utf-8"
    soup = BeautifulSoup(bytes(body), "lxml")
    for tag in soup(["script", "style", "noscript", "nav", "footer"]):
        tag.decompose()
    text = "\n".join(line.strip() for line in soup.get_text("\n").splitlines() if line.strip())
    if not text:
        raise ValueError("URL에서 읽을 본문을 찾지 못했습니다.")
    return [Document(page_content=text, metadata={"source": url, "page": "-"})]


def load_file(path: Path) -> list:
    from langchain_core.documents import Document

    suffix = path.suffix.lower()
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(str(path))
        documents = []
        for number, slide in enumerate(presentation.slides, 1):
            text = "\n".join(
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            )
            if text:
                documents.append(Document(text, metadata={"source": str(path), "page": number}))
        return documents
    if suffix == ".docx":
        from docx import Document as WordDocument

        document = WordDocument(str(path))
        text = "\n".join(paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip())
        return [Document(text, metadata={"source": str(path), "page": "-"})] if text else []
    if suffix == ".pdf":
        import pymupdf
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pdf = pymupdf.open(str(path))
        tessdata_dir = Path(os.environ.get("RAG_TESSDATA_DIR", Path(__file__).resolve().parents[1] / "data" / "tessdata"))
        document_title = extract_pdf_title(
            "\n".join(page.extract_text() or "" for page in reader.pages[:1]), path.stem
        )
        documents = []
        try:
            for number, page in enumerate(reader.pages, 1):
                text_layer = normalize_pdf_text(page.extract_text() or "")
                entries = structure_compact_faq_entries(text_layer, document_title)
                if not entries:
                    ocr_text = normalize_pdf_text(_ocr_pdf_page(pdf, number - 1, tessdata_dir))
                    entries = structure_pdf_entries(ocr_text, document_title)
                for entry in entries:
                    documents.append(Document(entry, metadata={"source": str(path), "page": number}))
        finally:
            pdf.close()
        if not documents:
            raise ValueError("PDF에서 읽을 수 있는 텍스트를 찾지 못했습니다.")
        return documents
    raise ValueError("PPTX와 DOCX 파일만 지원합니다.")
