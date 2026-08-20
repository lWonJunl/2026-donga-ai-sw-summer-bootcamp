from __future__ import annotations

import ipaddress
import socket
from pathlib import Path
from urllib.parse import urlparse

import requests


MAX_WEB_BYTES = 5 * 1024 * 1024


def validate_public_url(url: str) -> str:
    parsed = urlparse(url.strip())
    if parsed.scheme not in {"http", "https"} or not parsed.hostname:
        raise ValueError("http 또는 https URL만 사용할 수 있습니다.")
    if parsed.username or parsed.password:
        raise ValueError("인증정보가 포함된 URL은 사용할 수 없습니다.")
    try:
        addresses = {
            item[4][0] for item in socket.getaddrinfo(parsed.hostname, parsed.port or 443)
        }
    except socket.gaierror as error:
        raise ValueError("URL의 호스트를 확인할 수 없습니다.") from error
    for address in addresses:
        ip = ipaddress.ip_address(address)
        if not ip.is_global:
            raise ValueError("내부망·로컬 주소는 수집할 수 없습니다.")
    return parsed.geturl()


def load_url(url: str) -> list:
    from bs4 import BeautifulSoup
    from langchain_core.documents import Document

    url = validate_public_url(url)
    with requests.get(
        url,
        headers={"User-Agent": "EXAONE-Django-RAG/1.0"},
        timeout=(5, 20),
        stream=True,
        allow_redirects=False,
    ) as response:
        if 300 <= response.status_code < 400:
            raise ValueError("리디렉션 URL은 수집할 수 없습니다. 최종 URL을 입력하세요.")
        response.raise_for_status()
        content_type = response.headers.get("Content-Type", "").lower()
        if "text/html" not in content_type and "text/plain" not in content_type:
            raise ValueError("HTML 또는 텍스트 URL만 수집할 수 있습니다.")
        body = bytearray()
        for chunk in response.iter_content(64 * 1024):
            body.extend(chunk)
            if len(body) > MAX_WEB_BYTES:
                raise ValueError("URL 본문은 5MB 이하여야 합니다.")
    response.encoding = response.encoding or response.apparent_encoding or "utf-8"
    soup = BeautifulSoup(bytes(body), "lxml")
    for tag in soup(["script", "style", "noscript", "nav", "footer"]):
        tag.decompose()
    text = "\n".join(line.strip() for line in soup.get_text("\n").splitlines() if line.strip())
    if not text:
        raise ValueError("URL에서 읽을 본문을 찾지 못했습니다.")
    return [Document(page_content=text, metadata={"source": url, "page": "-"})]


def load_file(path: Path) -> list:
    from langchain_core.documents import Document

    suffix = path.suffix.lower()
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(str(path))
        documents = []
        for number, slide in enumerate(presentation.slides, 1):
            text = "\n".join(
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            )
            if text:
                documents.append(Document(text, metadata={"source": str(path), "page": number}))
        return documents
    if suffix == ".docx":
        from docx import Document as WordDocument

        document = WordDocument(str(path))
        text = "\n".join(paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip())
        return [Document(text, metadata={"source": str(path), "page": "-"})] if text else []
    raise ValueError("PPTX와 DOCX 파일만 지원합니다.")
